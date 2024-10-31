from pptx import Presentation

# Load the presentation
presentation_path = '/home/user/Desktop/33_1.pptx'
prs = Presentation(presentation_path)

# Access slide 4 (index 3 since it's zero-based)
slide = prs.slides[3]

# Locate the table on the slide
table = None
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        break

# Check if the table was found
if table is not None:
    # Modify the first row of the table
    table.cell(0, 0).text = "T1"
    table.cell(0, 1).text = "T2"
    table.cell(0, 2).text = "T3"
    table.cell(0, 3).text = "T4"
    
    # Save the changes
    prs.save(presentation_path)
    print("Python Script Run Successfully!!!")
else:
    print("No table found on slide 4.")
